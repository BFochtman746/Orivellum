"""Generator integrity tests — Excel, DOCX, PDF, PPTX.

Each test asserts:
  1. The output file exists on disk.
  2. The file is not zero-byte (no silent empty-file failure).
  3. The file can be opened by its native library (no structural corruption).
  4. A library doc was registered: GET-equivalent call returns a doc dict
     with readiness == "ready".

Scenarios:
  - Empty Work (no knowledge, no documents, no tasks) for all four formats.
  - Work title containing XML/HTML special chars (<, &, >) — PDF and DOCX
    generators use ReportLab markup and python-docx XML; escaping must be correct.
  - Large Work (500+ knowledge items) — exercises truncation paths.
  - Work with very long per-item text (>500 chars per item) — slice guards.
"""
from __future__ import annotations

import threading
from pathlib import Path
from typing import Any
from unittest.mock import patch, MagicMock

import pytest


# ── Fixtures ──────────────────────────────────────────────────────────────────

@pytest.fixture()
def db(tmp_path):
    """Return a real OrivellumDB backed by a temporary SQLite file."""
    from orivellum.database.db import OrivellumDB
    return OrivellumDB(str(tmp_path / "test.db"))


@pytest.fixture()
def cfg(tmp_path):
    """Return a minimal config stub pointing at the test tmpdir."""
    m = MagicMock()
    m.data_dir = str(tmp_path)
    m.llm = MagicMock()
    m.llm.model = "gpt-3.5-turbo"
    return m


def _make_work(db, title: str = "Test Work") -> dict:
    return db.create_work(title=title, description="A test work for generation.")


def _add_knowledge(db, work_id: str, n: int, text_template: str = "Fact number {i}.") -> None:
    """Insert *n* knowledge items into *work_id*."""
    for i in range(n):
        db.create_knowledge_item(
            work_id=work_id,
            kind="fact",
            text=text_template.format(i=i),
            subject=f"subject_{i}",
            predicate="is",
            obj=f"obj_{i}",
            confidence=0.9,
        )


def _patched_generate(fn, *args, **kwargs):
    """Run a generator with non-fatal mocks for network-dependent side-effects."""
    with (
        patch("orivellum.capabilities.embeddings.embed_chunks_for_doc", return_value=0),
        patch("orivellum.capabilities.persist.record_provenance", return_value=None),
    ):
        return fn(*args, **kwargs)


def _assert_valid_doc(db, doc_id: str) -> None:
    """Assert the registered library doc is fetchable and marked ready."""
    doc = db.get_document(doc_id)
    assert doc is not None, f"doc_id {doc_id!r} not found in library after generation"
    assert doc.get("readiness") == "ready", (
        f"Expected readiness='ready', got {doc.get('readiness')!r} for doc {doc_id}"
    )


def _assert_valid_file(fpath: Path, min_bytes: int = 1) -> None:
    """Assert the file exists and is non-empty."""
    assert fpath.exists(), f"Generator did not create file: {fpath}"
    size = fpath.stat().st_size
    assert size >= min_bytes, (
        f"Generator produced a {size}-byte file (suspected corruption): {fpath}"
    )


# ── Excel ─────────────────────────────────────────────────────────────────────

class TestExcelIntegrity:

    def test_empty_work(self, db, cfg):
        """Excel generation on a Work with no knowledge/docs/tasks must not fail."""
        import openpyxl
        from orivellum.capabilities.generate import generate_excel

        work = _make_work(db, "Empty Excel Work")
        fpath, doc_id = _patched_generate(generate_excel, work["id"], db, cfg)

        _assert_valid_file(fpath)
        _assert_valid_doc(db, doc_id)

        # File must be a valid xlsx (openpyxl can open it without error)
        wb = openpyxl.load_workbook(str(fpath))
        assert "Summary" in wb.sheetnames
        assert "Knowledge" in wb.sheetnames

    def test_large_work_500_knowledge_items(self, db, cfg):
        """Excel generation handles 500 knowledge items without truncation bugs."""
        import openpyxl
        from orivellum.capabilities.generate import generate_excel

        work = _make_work(db, "Large Knowledge Work")
        _add_knowledge(db, work["id"], 500)
        fpath, doc_id = _patched_generate(generate_excel, work["id"], db, cfg)

        _assert_valid_file(fpath, min_bytes=5_000)
        _assert_valid_doc(db, doc_id)

        wb = openpyxl.load_workbook(str(fpath))
        ws_kn = wb["Knowledge"]
        # Header row + up to 500 data rows
        data_rows = ws_kn.max_row - 1   # subtract header
        assert data_rows == 500, f"Expected 500 knowledge rows, got {data_rows}"

    def test_special_chars_in_work_title(self, db, cfg):
        """Excel workbook must be valid when the Work title contains &, <, >."""
        import openpyxl
        from orivellum.capabilities.generate import generate_excel

        work = _make_work(db, "Research & <Analysis> — 'Quotes' & \"More\"")
        _add_knowledge(db, work["id"], 3, text_template="Finding with <special> & chars {i}.")
        fpath, doc_id = _patched_generate(generate_excel, work["id"], db, cfg)

        _assert_valid_file(fpath)
        _assert_valid_doc(db, doc_id)
        # openpyxl would raise on a corrupt file
        wb = openpyxl.load_workbook(str(fpath))
        assert wb["Summary"]["A1"].value is not None


# ── DOCX ──────────────────────────────────────────────────────────────────────

class TestDocxIntegrity:

    def test_empty_work(self, db, cfg):
        """DOCX generation on a Work with no knowledge must produce a valid file."""
        from docx import Document
        from orivellum.capabilities.generate import generate_docx_report

        work = _make_work(db, "Empty DOCX Work")
        fpath, doc_id = _patched_generate(generate_docx_report, work["id"], db, cfg)

        _assert_valid_file(fpath)
        _assert_valid_doc(db, doc_id)
        # Document() would raise on structural corruption
        doc = Document(str(fpath))
        assert any("Research Report" in p.text or "Empty DOCX Work" in p.text
                   for p in doc.paragraphs)

    def test_very_long_knowledge_text(self, db, cfg):
        """DOCX generator must not crash when knowledge items exceed 500 chars."""
        from docx import Document
        from orivellum.capabilities.generate import generate_docx_report

        work = _make_work(db, "Long Text Work")
        long_text = "A" * 2_000   # well over the 500-char slice in generate_docx_report
        _add_knowledge(db, work["id"], 10, text_template=long_text)
        fpath, doc_id = _patched_generate(generate_docx_report, work["id"], db, cfg)

        _assert_valid_file(fpath)
        _assert_valid_doc(db, doc_id)
        Document(str(fpath))   # validates XML structure

    def test_special_chars_in_work_title(self, db, cfg):
        """DOCX generator must survive XML-special chars in title and description."""
        from docx import Document
        from orivellum.capabilities.generate import generate_docx_report

        # python-docx handles its own XML escaping; verify no lxml crash
        work = _make_work(db, "Work <Analysis> & 'Findings'")
        _add_knowledge(db, work["id"], 5, text_template="Item with <angle> & ampersand {i}")
        fpath, doc_id = _patched_generate(generate_docx_report, work["id"], db, cfg)

        _assert_valid_file(fpath)
        _assert_valid_doc(db, doc_id)
        Document(str(fpath))


# ── PDF ───────────────────────────────────────────────────────────────────────

class TestPdfIntegrity:

    def test_empty_work(self, db, cfg):
        """PDF generation on a Work with no knowledge must produce a valid PDF."""
        from orivellum.capabilities.generate import generate_pdf_report

        work = _make_work(db, "Empty PDF Work")
        fpath, doc_id = _patched_generate(generate_pdf_report, work["id"], db, cfg)

        _assert_valid_file(fpath)
        _assert_valid_doc(db, doc_id)
        # PDF magic bytes
        assert fpath.read_bytes()[:4] == b"%PDF", "Output is not a valid PDF (missing %PDF header)"

    def test_special_chars_escaped_in_pdf(self, db, cfg):
        """PDF generator must XML-escape & < > in titles and knowledge text."""
        from orivellum.capabilities.generate import generate_pdf_report

        work = _make_work(db, "Work & <Findings> — 'Study'")
        _add_knowledge(db, work["id"], 5, text_template="Finding: x < y & z > 0 {i}")
        fpath, doc_id = _patched_generate(generate_pdf_report, work["id"], db, cfg)

        _assert_valid_file(fpath)
        _assert_valid_doc(db, doc_id)
        # ReportLab raises on bad markup; reaching here means escaping worked
        assert fpath.read_bytes()[:4] == b"%PDF"

    def test_very_long_knowledge_text(self, db, cfg):
        """PDF generator must not crash when knowledge items exceed 400 chars."""
        from orivellum.capabilities.generate import generate_pdf_report

        work = _make_work(db, "PDF Long Text Work")
        long_text = "B" * 1_500
        _add_knowledge(db, work["id"], 15, text_template=long_text)
        fpath, doc_id = _patched_generate(generate_pdf_report, work["id"], db, cfg)

        _assert_valid_file(fpath)
        _assert_valid_doc(db, doc_id)
        assert fpath.read_bytes()[:4] == b"%PDF"

    def test_large_work_500_knowledge_items(self, db, cfg):
        """PDF generator handles 500 knowledge items without error."""
        from orivellum.capabilities.generate import generate_pdf_report

        work = _make_work(db, "PDF Large Work")
        _add_knowledge(db, work["id"], 500)
        fpath, doc_id = _patched_generate(generate_pdf_report, work["id"], db, cfg)

        # 500 short-text items: file is well-formed but compact; 2 000 B is generous
        _assert_valid_file(fpath, min_bytes=2_000)
        _assert_valid_doc(db, doc_id)
        assert fpath.read_bytes()[:4] == b"%PDF"


# ── PPTX ──────────────────────────────────────────────────────────────────────

class TestPptxIntegrity:

    def test_empty_work(self, db, cfg):
        """PPTX generation on a Work with no knowledge must produce a valid deck."""
        from pptx import Presentation
        from orivellum.capabilities.generate import generate_pptx

        work = _make_work(db, "Empty Slides Work")
        fpath, doc_id = _patched_generate(generate_pptx, work["id"], db, cfg)

        _assert_valid_file(fpath)
        _assert_valid_doc(db, doc_id)
        prs = Presentation(str(fpath))
        # Title slide + overview slide even when empty
        assert len(prs.slides) >= 2

    def test_special_chars_in_work_title(self, db, cfg):
        """PPTX generator must not crash with XML-special chars in title."""
        from pptx import Presentation
        from orivellum.capabilities.generate import generate_pptx

        work = _make_work(db, "Slides & <Research> — \"Study\"")
        _add_knowledge(db, work["id"], 5, text_template="Bullet: a < b & c > d {i}")
        fpath, doc_id = _patched_generate(generate_pptx, work["id"], db, cfg)

        _assert_valid_file(fpath)
        _assert_valid_doc(db, doc_id)
        Presentation(str(fpath))

    def test_large_work_500_knowledge_items(self, db, cfg):
        """PPTX generator handles 500 knowledge items (multiple kind-slides)."""
        from pptx import Presentation
        from orivellum.capabilities.generate import generate_pptx

        work = _make_work(db, "PPTX Large Work")
        _add_knowledge(db, work["id"], 500)
        fpath, doc_id = _patched_generate(generate_pptx, work["id"], db, cfg)

        _assert_valid_file(fpath, min_bytes=5_000)
        _assert_valid_doc(db, doc_id)
        prs = Presentation(str(fpath))
        # Title + overview + at least one knowledge slide
        assert len(prs.slides) >= 3

    def test_very_long_knowledge_text(self, db, cfg):
        """PPTX generator slices long text; deck must still be valid."""
        from pptx import Presentation
        from orivellum.capabilities.generate import generate_pptx

        work = _make_work(db, "PPTX Long Text Work")
        long_text = "C" * 1_000
        _add_knowledge(db, work["id"], 20, text_template=long_text)
        fpath, doc_id = _patched_generate(generate_pptx, work["id"], db, cfg)

        _assert_valid_file(fpath)
        _assert_valid_doc(db, doc_id)
        Presentation(str(fpath))


# ── Registration atomicity ─────────────────────────────────────────────────────

class TestRegistrationAtomicity:

    def test_doc_id_returned_is_fetchable(self, db, cfg):
        """_register_output must return a doc_id that get_document() can find."""
        from orivellum.capabilities.generate import generate_excel

        work = _make_work(db, "Registration Test")
        _add_knowledge(db, work["id"], 2)
        fpath, doc_id = _patched_generate(generate_excel, work["id"], db, cfg)

        assert doc_id and isinstance(doc_id, str), "doc_id must be a non-empty string"
        doc = db.get_document(doc_id)
        assert doc is not None
        assert doc["readiness"] == "ready"
        # Kind must match the file extension
        assert doc["kind"] == "xlsx"
        # Tier must be 'artifact' (generation outputs must not pollute canon corpus)
        assert doc.get("tier") == "artifact", (
            f"Generated doc must be ARTIFACT tier, got {doc.get('tier')!r}"
        )

    def test_unscoped_work_id_registered(self, db, cfg):
        """Generated doc must still be registered even when work has no knowledge."""
        from orivellum.capabilities.generate import generate_pptx

        work = _make_work(db, "Minimal Work")
        # Zero knowledge items — exercises the empty-text fallback in _register_output
        fpath, doc_id = _patched_generate(generate_pptx, work["id"], db, cfg)

        doc = db.get_document(doc_id)
        assert doc is not None
        assert doc["readiness"] == "ready"
        assert doc["work_id"] == work["id"]
