"""Audit remediation — Workshop sandbox, backup restore, Excel polish.

Acceptance criteria from the remediation packet:
- Workshop script execution is sandboxed: a script that deliberately attempts
  an outbound network request must FAIL, and the child environment must not
  contain the parent's secrets.
- Backup → restore round-trip works end-to-end: stage a restore, apply it,
  and confirm the database file was swapped and a safety snapshot was kept.
- generate_excel output contains at least one chart and a frozen header row.
"""
from __future__ import annotations

import os
import tempfile
import unittest
import zipfile
from pathlib import Path

from fastapi.testclient import TestClient

from orivellum.database.db import OrivellumDB
from tests.conftest import AUTH_HEADERS


def _make_app(tmp: str):
    from orivellum.api import _deps
    from orivellum.api.app import app
    from orivellum.configuration.config import OrivellumConfig

    cfg = OrivellumConfig(data_dir=tmp)
    db = OrivellumDB(str(Path(tmp) / "test.db"))
    _deps.init(db=db, cfg=cfg)
    return app, db, cfg


# ── Workshop sandbox ─────────────────────────────────────────────────────────

class TestWorkshopSandbox(unittest.TestCase):
    """The no-network rule must be enforced, not just requested in the prompt."""

    def setUp(self):
        self._tmp = tempfile.TemporaryDirectory()
        _, self.db, self.cfg = _make_app(self._tmp.name)

    def tearDown(self):
        self._tmp.cleanup()

    def _run(self, script: str) -> dict:
        from orivellum.capabilities.workshop import _run_script_safely
        return _run_script_safely(
            script, output_path=str(Path(self._tmp.name) / "out.bin"),
            max_retries=0, cfg=self.cfg, db=self.db, request="test",
        )

    def test_outbound_network_request_fails(self):
        script = (
            "import urllib.request\n"
            "urllib.request.urlopen('http://127.0.0.1:1/x', timeout=2)\n"
        )
        result = self._run(script)
        self.assertFalse(result["ok"])
        self.assertIn("disabled", (result.get("stderr") or "") + result.get("error", ""))

    def test_raw_socket_connection_fails(self):
        result = self._run(
            "import socket\n"
            "socket.socket(socket.AF_INET, socket.SOCK_STREAM)\n"
        )
        self.assertFalse(result["ok"])
        self.assertIn("disabled", (result.get("stderr") or "") + result.get("error", ""))

    def test_low_level_socket_module_also_denied(self):
        result = self._run("import _socket\n_socket.socket()\n")
        self.assertFalse(result["ok"])

    def test_reportlab_pdf_generation_works_in_sandbox(self):
        """reportlab imports urllib internals — must still work offline."""
        out = Path(self._tmp.name) / "out.pdf"
        result = self._run(
            "from reportlab.pdfgen import canvas\n"
            f"c = canvas.Canvas({str(out)!r})\n"
            "c.drawString(100, 750, 'sandbox pdf test')\n"
            "c.save()\n"
        )
        self.assertTrue(result["ok"], result.get("error") or result.get("stderr"))
        self.assertTrue(out.exists() and out.stat().st_size > 100)

    def test_pptx_generation_works_in_sandbox(self):
        out = Path(self._tmp.name) / "out.pptx"
        result = self._run(
            "from pptx import Presentation\n"
            "p = Presentation()\n"
            "p.slides.add_slide(p.slide_layouts[6])\n"
            f"p.save({str(out)!r})\n"
        )
        self.assertTrue(result["ok"], result.get("error") or result.get("stderr"))
        self.assertTrue(out.exists() and out.stat().st_size > 1000)

    def test_parent_secrets_not_in_child_env(self):
        os.environ["ORIVELLUM_TEST_FAKE_SECRET"] = "leak-canary"
        try:
            result = self._run(
                "import os\n"
                "assert 'ORIVELLUM_TEST_FAKE_SECRET' not in os.environ, 'secret leaked'\n"
                "assert 'SESSION_SECRET' not in os.environ\n"
                "assert 'TAVILY_API_KEY' not in os.environ\n"
                "print('env clean')\n"
            )
        finally:
            del os.environ["ORIVELLUM_TEST_FAKE_SECRET"]
        self.assertTrue(result["ok"], result.get("error") or result.get("stderr"))
        self.assertIn("env clean", result["stdout"])

    def test_legitimate_document_script_still_works(self):
        out = Path(self._tmp.name) / "out.txt"
        result = self._run(f"open({str(out)!r}, 'w').write('hello document')\n")
        self.assertTrue(result["ok"], result.get("error") or result.get("stderr"))
        self.assertEqual(out.read_text(), "hello document")


# ── Backup restore ───────────────────────────────────────────────────────────

class TestBackupRestore(unittest.TestCase):
    def setUp(self):
        self._tmp = tempfile.TemporaryDirectory()
        self.app, self.db, self.cfg = _make_app(self._tmp.name)
        self.client = TestClient(self.app)

    def tearDown(self):
        self._tmp.cleanup()

    def _make_backup_zip(self, name: str = "orivellum_backup_test.zip") -> Path:
        bd = Path(self.cfg.data_dir) / "backups"
        bd.mkdir(parents=True, exist_ok=True)
        zp = bd / name
        with zipfile.ZipFile(zp, "w") as zf:
            zf.writestr("orivellum.db", b"RESTORED-DB-CONTENT")
            zf.writestr("library/doc.txt", b"restored library file")
        return zp

    def test_stage_restore_and_cancel(self):
        self._make_backup_zip()
        r = self.client.post("/api/backups/orivellum_backup_test.zip/restore",
                             headers=AUTH_HEADERS)
        self.assertEqual(r.status_code, 200, r.text)
        self.assertTrue(r.json()["staged"])
        self.assertTrue((Path(self.cfg.data_dir) / "restore-pending.zip").exists())

        r = self.client.get("/api/backups/restore/pending", headers=AUTH_HEADERS)
        self.assertTrue(r.json()["pending"])
        r = self.client.delete("/api/backups/restore/pending", headers=AUTH_HEADERS)
        self.assertTrue(r.json()["cancelled"])
        self.assertFalse((Path(self.cfg.data_dir) / "restore-pending.zip").exists())

    def test_stage_rejects_archive_without_db(self):
        bd = Path(self.cfg.data_dir) / "backups"
        bd.mkdir(parents=True, exist_ok=True)
        with zipfile.ZipFile(bd / "bad.zip", "w") as zf:
            zf.writestr("random.txt", b"nope")
        r = self.client.post("/api/backups/bad.zip/restore", headers=AUTH_HEADERS)
        self.assertEqual(r.status_code, 422)

    def test_apply_pending_restore_round_trip(self):
        """The startup hook swaps the DB, keeps a safety snapshot."""
        from orivellum.api.app import _apply_pending_restore

        data_dir = Path(self.cfg.data_dir)
        db_path = Path(self.cfg.db_path)
        db_path.parent.mkdir(parents=True, exist_ok=True)
        db_path.write_bytes(b"CURRENT-DB-CONTENT")

        zp = self._make_backup_zip()
        import shutil
        shutil.copy2(zp, data_dir / "restore-pending.zip")

        _apply_pending_restore(self.cfg)

        # DB swapped in
        self.assertEqual((data_dir / "orivellum.db").read_bytes(), b"RESTORED-DB-CONTENT")
        # Library restored
        self.assertEqual((data_dir / "library" / "doc.txt").read_bytes(),
                         b"restored library file")
        # Pending marker consumed
        self.assertFalse((data_dir / "restore-pending.zip").exists())
        # Safety snapshot holds the old DB
        safety = list(data_dir.glob("restore-safety-*"))
        self.assertEqual(len(safety), 1)
        self.assertEqual((safety[0] / "orivellum.db").read_bytes(), b"CURRENT-DB-CONTENT")

    def test_apply_bad_archive_never_bricks_startup(self):
        from orivellum.api.app import _apply_pending_restore

        data_dir = Path(self.cfg.data_dir)
        db_path = Path(self.cfg.db_path)
        db_path.parent.mkdir(parents=True, exist_ok=True)
        db_path.write_bytes(b"CURRENT-DB-CONTENT")
        (data_dir / "restore-pending.zip").write_bytes(b"not a zip at all")

        _apply_pending_restore(self.cfg)  # must not raise

        self.assertEqual((data_dir / "orivellum.db").read_bytes(), b"CURRENT-DB-CONTENT")
        self.assertFalse((data_dir / "restore-pending.zip").exists())
        self.assertTrue((data_dir / "restore-failed.zip").exists())


# ── Excel polish ─────────────────────────────────────────────────────────────

class TestExcelPolish(unittest.TestCase):
    def test_generated_workbook_has_chart_freeze_and_filter(self):
        import openpyxl

        from orivellum.capabilities.generate import generate_excel

        with tempfile.TemporaryDirectory() as tmp:
            _, db, cfg = _make_app(tmp)
            work = db.create_work("Excel Test Work", "for chart checks")
            db.create_knowledge_item(work_id=work["id"], kind="fact",
                                     text="test fact", confidence=0.9)
            db.create_task(work["id"], "a task")

            fpath, _doc_id = generate_excel(work["id"], db, cfg)
            wb = openpyxl.load_workbook(str(fpath))

            self.assertGreaterEqual(len(wb["Summary"]._charts), 1,
                                    "Summary sheet must contain a chart")
            self.assertEqual(wb["Knowledge"].freeze_panes, "A2")
            self.assertEqual(wb["Tasks"].freeze_panes, "A2")
            self.assertTrue(wb["Knowledge"].auto_filter.ref,
                            "Knowledge sheet must have an autofilter")


if __name__ == "__main__":
    unittest.main()
