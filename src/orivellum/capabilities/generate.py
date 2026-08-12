"""Document generation capabilities — Excel, DOCX, PDF, PPTX.

Each generator:
  1. Pulls a Work's knowledge, documents, and tasks from the database.
  2. Produces a file in data/outputs/generate/{work_id}/
  3. Registers the output as a library document (tier=artifact) with
     provenance metadata so it is tracked, searchable, and recallable.
  4. Returns (file_path, doc_id) for the caller to serve.

No-pollution guarantee: generated outputs are always registered as
ARTIFACT tier — they are never promoted into a Work's canonical corpus.
"""

from __future__ import annotations

import hashlib
import logging
import zipfile
from datetime import UTC, datetime
from pathlib import Path
from typing import TYPE_CHECKING

from orivellum.version import code_version as _code_version

if TYPE_CHECKING:
    from orivellum.configuration.config import OrivellumConfig
    from orivellum.database.db import OrivellumDB

logger = logging.getLogger("orivellum.generate")

# ── Shared helpers ─────────────────────────────────────────────────────────────


def _now_label() -> str:
    return datetime.now(UTC).strftime("%Y%m%d_%H%M")


def _ensure_dir(cfg: OrivellumConfig, work_id: str) -> Path:
    """Return (and create) data/outputs/generate/{work_id}/"""
    d = Path(cfg.data_dir) / "outputs" / "generate" / work_id
    d.mkdir(parents=True, exist_ok=True)
    return d


_CHUNK_SIZE = 1_000  # characters per FTS chunk


def _register_output(
    doc_path: Path,
    work_id: str | None,
    db: OrivellumDB,
    cfg: OrivellumConfig,
    format_label: str,
    title: str,
    text_content: str = "",
) -> str:
    """Register a generated file as a library document; return its doc_id.

    work_id may be None for unscoped outputs (e.g. tax packages that span all Works).

    Amendment-1 save/process/recall invariant: the full text is written to
    documents.extracted_text AND split into chunks that are inserted into
    chunks + chunks_fts so the document is findable via the application's
    document full-text search path (which queries chunks_fts, not extracted_text).
    """
    # Amendment-1 path fix: content_path must be relative to lib_root (data_dir/library)
    # because every Library endpoint resolves paths as lib_root / content_path.
    # _ensure_lib_symlink creates a zero-cost symlink under lib_root/generated/ for
    # files stored outside lib_root (outputs/generate/...) and returns the correct
    # lib-root-relative path so Library downloads, reprocess, and resolution all work.
    from orivellum.capabilities.persist import _ensure_lib_symlink

    lib_root = Path(cfg.data_dir) / "library"
    rel = _ensure_lib_symlink(doc_path, lib_root)
    sha = hashlib.sha256(doc_path.read_bytes()).hexdigest() if doc_path.exists() else None
    scope_label = work_id or "library"

    doc = db.create_document(
        title=title,
        source=f"generated/{format_label}/{scope_label}",
        sha256=sha,
        kind=doc_path.suffix.lstrip("."),
        work_id=work_id,
        content_path=rel,
        meta={
            "provenance": "generation",
            "origin_id": work_id,
            "format": format_label,
            "generated_at": datetime.now(UTC).isoformat(),
            "code_version": _code_version(),
        },
        tier="artifact",
        doc_type="generated",
        doc_type_by="rule:system-output",
    )
    doc_id = doc["id"]

    body = text_content or f"Generated {format_label} output for work {work_id}"

    # Mark the document ready and store the full text for direct lookup
    db.update_document_extracted(
        doc_id,
        extracted_text=body,
        word_count=max(len(body.split()), 1),
        readiness="ready",
    )

    # Split into overlapping chunks and index into chunks_fts so the document
    # is searchable via the application's standard document search path.
    words = body.split()
    # Build ~1 000-char chunks from the word list (no mid-word splits)
    chunk_texts: list[str] = []
    current: list[str] = []
    current_len = 0
    for word in words:
        wl = len(word) + 1
        if current_len + wl > _CHUNK_SIZE and current:
            chunk_texts.append(" ".join(current))
            current = []
            current_len = 0
        current.append(word)
        current_len += wl
    if current:
        chunk_texts.append(" ".join(current))

    for page, chunk in enumerate(chunk_texts):
        try:
            db.add_chunk(doc_id, chunk, page=page)
        except Exception:
            logger.warning("_register_output: failed to add chunk %d for doc %s", page, doc_id)

    # Amendment-1: embed all chunks in the background so generated documents
    # are semantically searchable without waiting for the nightly backfill.
    # Wrapped in a try/except so thread failures are logged, never raised.
    import threading as _t

    from orivellum.capabilities.persist import record_provenance as _prov

    def _embed_bg(_doc_id=doc_id, _db=db) -> None:
        try:
            from orivellum.capabilities.embeddings import embed_chunks_for_doc

            embed_chunks_for_doc(_doc_id, _db)
        except Exception as _exc:
            logger.debug("_register_output background embed failed (non-fatal): %s", _exc)

    _t.Thread(target=_embed_bg, daemon=True).start()
    _prov(doc_id, "generation", db, work_id=work_id)

    return doc_id


def _format_knowledge_text(items: list[dict], max_items: int = 500) -> str:
    """Build a plain-text version of knowledge items for embedding."""
    lines = []
    for it in items[:max_items]:
        kind = it.get("kind", "note")
        text = it.get("text", "").strip()
        if text:
            lines.append(f"[{kind.upper()}] {text}")
    return "\n".join(lines)


def _slug(text: str, max_len: int = 40) -> str:
    import re

    s = re.sub(r"[^a-z0-9]+", "_", text.lower().strip())
    return s[:max_len].strip("_")


# ── Excel ──────────────────────────────────────────────────────────────────────


def generate_excel(work_id: str, db: OrivellumDB, cfg: OrivellumConfig) -> tuple[Path, str]:
    """Generate an xlsx workbook summarising a Work; return (file_path, doc_id)."""
    import openpyxl
    from openpyxl.styles import Alignment, Border, Font, PatternFill, Side
    from openpyxl.utils import get_column_letter

    work = db.get_work(work_id)
    if not work:
        raise ValueError(f"Work {work_id!r} not found")

    knowledge = db.list_knowledge(work_id=work_id, limit=500)
    tasks = db.list_tasks(work_id=work_id)
    # Documents via DB
    with db._lock:
        doc_rows = db._conn.execute(
            "SELECT id, title, kind, readiness, source, created_at FROM documents WHERE work_id=? ORDER BY created_at DESC",
            (work_id,),
        ).fetchall()
    doc_list = [dict(r) for r in doc_rows]

    wb = openpyxl.Workbook()

    # ── Styles ──
    hdr_fill = PatternFill("solid", fgColor="1E293B")
    hdr_font = Font(bold=True, color="FFFFFF", size=10)
    title_font = Font(bold=True, size=14)
    sub_font = Font(size=10, color="64748B")
    thin = Side(style="thin", color="E2E8F0")
    border = Border(left=thin, right=thin, top=thin, bottom=thin)

    def _apply_header(ws, headers: list[str], row: int = 1) -> None:
        for col, h in enumerate(headers, 1):
            cell = ws.cell(row=row, column=col, value=h)
            cell.font = hdr_font
            cell.fill = hdr_fill
            cell.alignment = Alignment(horizontal="center", vertical="center")

    def _autowidth(ws) -> None:
        for col in ws.columns:
            max_len = max((len(str(c.value or "")) for c in col), default=8)
            ws.column_dimensions[get_column_letter(col[0].column)].width = min(max_len + 4, 60)

    # ── Sheet 1: Summary ──
    ws = wb.active
    ws.title = "Summary"
    ws["A1"] = work.get("title", "Untitled Work")
    ws["A1"].font = title_font
    ws["A2"] = work.get("description") or ""
    ws["A2"].font = sub_font
    ws["A4"] = "Generated"
    ws["B4"] = datetime.now(UTC).strftime("%Y-%m-%d %H:%M UTC")
    ws["A5"] = "Documents"
    ws["B5"] = len(doc_list)
    ws["A6"] = "Knowledge items"
    ws["B6"] = len(knowledge)
    ws["A7"] = "Open tasks"
    ws["B7"] = sum(1 for t in tasks if t.get("status") == "pending")
    for r in range(4, 8):
        ws[f"A{r}"].font = Font(bold=True, size=10)
    ws.column_dimensions["A"].width = 22
    ws.column_dimensions["B"].width = 40

    # ── Summary chart: contents at a glance ──
    from openpyxl.chart import BarChart, Reference

    chart = BarChart()
    chart.type = "col"
    chart.title = "Work contents"
    chart.legend = None
    chart.y_axis.title = "Count"
    chart.add_data(Reference(ws, min_col=2, min_row=5, max_row=7))
    chart.set_categories(Reference(ws, min_col=1, min_row=5, max_row=7))
    chart.height, chart.width = 7, 12
    ws.add_chart(chart, "D4")

    # ── Sheet 2: Knowledge ──
    ws2 = wb.create_sheet("Knowledge")
    headers2 = ["#", "Kind", "Text", "Confidence", "Created"]
    _apply_header(ws2, headers2)
    for i, item in enumerate(knowledge, 1):
        text = (item.get("text") or "").replace("\n", " ").strip()[:300]
        ws2.append(
            [
                i,
                item.get("kind", "note"),
                text,
                round(item.get("confidence", 0) * 100) if item.get("confidence") else "",
                (item.get("created_at") or "")[:10],
            ]
        )
    for row in ws2.iter_rows(min_row=2):
        for cell in row:
            cell.border = border
    _autowidth(ws2)
    ws2.freeze_panes = "A2"
    if knowledge:
        ws2.auto_filter.ref = f"A1:{get_column_letter(len(headers2))}{len(knowledge) + 1}"

    # ── Sheet 3: Documents ──
    ws3 = wb.create_sheet("Documents")
    headers3 = ["#", "Title", "Kind", "Readiness", "Created"]
    _apply_header(ws3, headers3)
    for i, d in enumerate(doc_list, 1):
        ws3.append(
            [
                i,
                (d.get("title") or "")[:80],
                d.get("kind", ""),
                d.get("readiness", ""),
                (d.get("created_at") or "")[:10],
            ]
        )
    for row in ws3.iter_rows(min_row=2):
        for cell in row:
            cell.border = border
    _autowidth(ws3)
    ws3.freeze_panes = "A2"
    if doc_list:
        ws3.auto_filter.ref = f"A1:{get_column_letter(len(headers3))}{len(doc_list) + 1}"

    # ── Sheet 4: Tasks ──
    ws4 = wb.create_sheet("Tasks")
    headers4 = ["#", "Status", "Priority", "Text", "Created"]
    _apply_header(ws4, headers4)
    for i, t in enumerate(tasks, 1):
        ws4.append(
            [
                i,
                t.get("status", ""),
                t.get("priority", 0),
                (t.get("text") or "")[:200],
                (t.get("created_at") or "")[:10],
            ]
        )
    for row in ws4.iter_rows(min_row=2):
        for cell in row:
            cell.border = border
    _autowidth(ws4)
    ws4.freeze_panes = "A2"
    if tasks:
        ws4.auto_filter.ref = f"A1:{get_column_letter(len(headers4))}{len(tasks) + 1}"

    # ── Save ──
    out_dir = _ensure_dir(cfg, work_id)
    slug = _slug(work.get("title", "work"))
    fname = f"{slug}_{_now_label()}.xlsx"
    fpath = out_dir / fname
    wb.save(str(fpath))

    text = _format_knowledge_text(knowledge)
    title_out = f"Research Workbook — {work.get('title', 'Work')}"
    doc_id = _register_output(fpath, work_id, db, cfg, "excel", title_out, text)
    logger.info("Generated Excel: %s → doc %s", fpath.name, doc_id)
    return fpath, doc_id


# ── DOCX ───────────────────────────────────────────────────────────────────────


def generate_docx_report(work_id: str, db: OrivellumDB, cfg: OrivellumConfig) -> tuple[Path, str]:
    """Generate a .docx research report from a Work; return (file_path, doc_id)."""
    from docx import Document
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.shared import Pt, RGBColor

    work = db.get_work(work_id)
    if not work:
        raise ValueError(f"Work {work_id!r} not found")

    knowledge = db.list_knowledge(work_id=work_id, limit=500)
    tasks = db.list_tasks(work_id=work_id)
    with db._lock:
        doc_rows = db._conn.execute(
            "SELECT title, kind, readiness FROM documents WHERE work_id=? ORDER BY created_at DESC",
            (work_id,),
        ).fetchall()
    doc_list = [dict(r) for r in doc_rows]

    document = Document()

    # ── Title ──
    h1 = document.add_heading(work.get("title", "Research Report"), 0)
    h1.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = h1.runs[0]
    run.font.color.rgb = RGBColor(0x0F, 0x17, 0x2A)

    desc = work.get("description") or ""
    if desc:
        p = document.add_paragraph(desc)
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p.runs[0].font.color.rgb = RGBColor(0x64, 0x74, 0x8B)
        p.runs[0].font.italic = True

    meta_p = document.add_paragraph(f"Generated {datetime.now(UTC).strftime('%B %d, %Y')}")
    meta_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    meta_p.runs[0].font.size = Pt(9)
    meta_p.runs[0].font.color.rgb = RGBColor(0x94, 0xA3, 0xB8)
    document.add_page_break()

    # ── Stats ──
    document.add_heading("Overview", 1)
    stats_table = document.add_table(rows=1, cols=3)
    stats_table.style = "Table Grid"
    hdr_cells = stats_table.rows[0].cells
    hdr_cells[0].text = "Documents"
    hdr_cells[1].text = "Knowledge Items"
    hdr_cells[2].text = "Open Tasks"
    row_cells = stats_table.add_row().cells
    row_cells[0].text = str(len(doc_list))
    row_cells[1].text = str(len(knowledge))
    row_cells[2].text = str(sum(1 for t in tasks if t.get("status") == "pending"))
    document.add_paragraph("")

    # ── Knowledge by kind ──
    from collections import defaultdict

    by_kind: dict[str, list[dict]] = defaultdict(list)
    for item in knowledge:
        by_kind[item.get("kind", "note")].append(item)

    kind_order = ["fact", "claim", "note", "question", "action", "entity", "summary"]
    all_kinds = kind_order + [k for k in by_kind if k not in kind_order]

    for kind in all_kinds:
        items = by_kind.get(kind, [])
        if not items:
            continue
        document.add_heading(kind.title() + "s", 2)
        for item in items[:80]:
            text = (item.get("text") or "").strip()
            if text:
                p = document.add_paragraph(style="List Bullet")
                p.add_run(text[:500])

    # ── Documents ──
    document.add_heading("Source Documents", 1)
    if doc_list:
        t = document.add_table(rows=1, cols=3)
        t.style = "Table Grid"
        hc = t.rows[0].cells
        hc[0].text = "Title"
        hc[1].text = "Kind"
        hc[2].text = "Status"
        for d in doc_list[:50]:
            rc = t.add_row().cells
            rc[0].text = (d.get("title") or "")[:80]
            rc[1].text = d.get("kind", "")
            rc[2].text = d.get("readiness", "")
    document.add_paragraph("")

    # ── Tasks ──
    pending = [t for t in tasks if t.get("status") == "pending"]
    if pending:
        document.add_heading("Open Tasks", 1)
        for task in pending[:50]:
            p = document.add_paragraph(style="List Bullet")
            p.add_run((task.get("text") or "")[:200])

    # ── Save ──
    out_dir = _ensure_dir(cfg, work_id)
    slug = _slug(work.get("title", "work"))
    fname = f"{slug}_{_now_label()}.docx"
    fpath = out_dir / fname
    document.save(str(fpath))

    text = _format_knowledge_text(knowledge)
    title_out = f"Research Report — {work.get('title', 'Work')}"
    doc_id = _register_output(fpath, work_id, db, cfg, "docx", title_out, text)
    logger.info("Generated DOCX: %s → doc %s", fpath.name, doc_id)
    return fpath, doc_id


# ── PDF ────────────────────────────────────────────────────────────────────────


def generate_pdf_report(work_id: str, db: OrivellumDB, cfg: OrivellumConfig) -> tuple[Path, str]:
    """Generate a PDF research report from a Work; return (file_path, doc_id)."""
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
    from reportlab.lib.units import cm
    from reportlab.platypus import (
        HRFlowable,
        PageBreak,
        Paragraph,
        SimpleDocTemplate,
        Spacer,
        Table,
        TableStyle,
    )

    work = db.get_work(work_id)
    if not work:
        raise ValueError(f"Work {work_id!r} not found")

    knowledge = db.list_knowledge(work_id=work_id, limit=500)
    tasks = db.list_tasks(work_id=work_id)
    with db._lock:
        doc_rows = db._conn.execute(
            "SELECT title, kind, readiness FROM documents WHERE work_id=? ORDER BY created_at DESC",
            (work_id,),
        ).fetchall()
    doc_list = [dict(r) for r in doc_rows]

    out_dir = _ensure_dir(cfg, work_id)
    slug = _slug(work.get("title", "work"))
    fname = f"{slug}_{_now_label()}.pdf"
    fpath = out_dir / fname

    doc = SimpleDocTemplate(
        str(fpath),
        pagesize=A4,
        rightMargin=2 * cm,
        leftMargin=2 * cm,
        topMargin=2 * cm,
        bottomMargin=2 * cm,
    )

    styles = getSampleStyleSheet()
    title_style = ParagraphStyle(
        "CustomTitle",
        parent=styles["Title"],
        fontSize=28,
        spaceAfter=12,
        textColor=colors.HexColor("#0F172A"),
    )
    h1_style = ParagraphStyle(
        "H1",
        parent=styles["Heading1"],
        fontSize=16,
        spaceBefore=18,
        spaceAfter=8,
        textColor=colors.HexColor("#1E293B"),
        borderPad=4,
    )
    h2_style = ParagraphStyle(
        "H2",
        parent=styles["Heading2"],
        fontSize=13,
        spaceBefore=14,
        spaceAfter=6,
        textColor=colors.HexColor("#334155"),
    )
    ParagraphStyle(
        "Body",
        parent=styles["Normal"],
        fontSize=10,
        leading=15,
        spaceAfter=4,
        textColor=colors.HexColor("#374151"),
    )
    bullet_style = ParagraphStyle(
        "Bullet",
        parent=styles["Normal"],
        fontSize=10,
        leading=14,
        leftIndent=20,
        bulletIndent=8,
        spaceAfter=3,
        textColor=colors.HexColor("#374151"),
    )
    meta_style = ParagraphStyle(
        "Meta",
        parent=styles["Normal"],
        fontSize=9,
        textColor=colors.HexColor("#94A3B8"),
        spaceAfter=6,
    )

    def _xml(s: str) -> str:
        """XML-escape a user-controlled string for use inside ReportLab Paragraph markup."""
        return s.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")

    story = []

    # Title — must be XML-escaped (title/description are user-controlled)
    title = _xml(work.get("title", "Research Report"))
    story.append(Paragraph(title, title_style))
    desc = work.get("description") or ""
    if desc:
        story.append(Paragraph(f"<i>{_xml(desc[:300])}</i>", meta_style))
    story.append(
        Paragraph(
            f"Generated {datetime.now(UTC).strftime('%B %d, %Y')}",
            meta_style,
        )
    )
    story.append(HRFlowable(width="100%", thickness=1, color=colors.HexColor("#E2E8F0")))
    story.append(Spacer(1, 0.5 * cm))

    # Stats table
    story.append(Paragraph("Overview", h1_style))
    tbl_data = [
        ["Documents", "Knowledge Items", "Open Tasks"],
        [
            str(len(doc_list)),
            str(len(knowledge)),
            str(sum(1 for t in tasks if t.get("status") == "pending")),
        ],
    ]
    tbl = Table(tbl_data, colWidths=[5 * cm, 5 * cm, 5 * cm])
    tbl.setStyle(
        TableStyle(
            [
                ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#1E293B")),
                ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                ("FONTSIZE", (0, 0), (-1, -1), 10),
                ("ALIGN", (0, 0), (-1, -1), "CENTER"),
                ("GRID", (0, 0), (-1, -1), 0.5, colors.HexColor("#E2E8F0")),
                ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#F8FAFC")]),
                ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
                ("TOPPADDING", (0, 0), (-1, -1), 6),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 6),
            ]
        )
    )
    story.append(tbl)
    story.append(Spacer(1, 0.5 * cm))

    # Knowledge by kind
    story.append(Paragraph("Knowledge Base", h1_style))
    from collections import defaultdict

    by_kind: dict[str, list[dict]] = defaultdict(list)
    for item in knowledge:
        by_kind[item.get("kind", "note")].append(item)

    kind_order = ["fact", "claim", "note", "question", "action", "entity", "summary"]
    all_kinds = kind_order + [k for k in by_kind if k not in kind_order]

    for kind in all_kinds:
        items = by_kind.get(kind, [])
        if not items:
            continue
        story.append(Paragraph(f"{kind.title()}s ({len(items)})", h2_style))
        for item in items[:60]:
            text = (item.get("text") or "").strip()[:400]
            if text:
                # Escape XML special chars
                text = text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
                story.append(Paragraph(f"• {text}", bullet_style))

    # Documents — all cell values through _xml() because ReportLab can
    # interpret strings containing &/</> as markup even in Table cells.
    if doc_list:
        story.append(PageBreak())
        story.append(Paragraph("Source Documents", h1_style))
        doc_data = [["Title", "Kind", "Status"]]
        for d in doc_list[:40]:
            doc_data.append(
                [
                    _xml((d.get("title") or "")[:60]),
                    _xml(d.get("kind", "")),
                    _xml(d.get("readiness", "")),
                ]
            )
        doc_tbl = Table(doc_data, colWidths=[10 * cm, 3 * cm, 3 * cm])
        doc_tbl.setStyle(
            TableStyle(
                [
                    ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#1E293B")),
                    ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                    ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                    ("FONTSIZE", (0, 0), (-1, -1), 9),
                    ("GRID", (0, 0), (-1, -1), 0.5, colors.HexColor("#E2E8F0")),
                    (
                        "ROWBACKGROUNDS",
                        (0, 1),
                        (-1, -1),
                        [colors.white, colors.HexColor("#F8FAFC")],
                    ),
                    ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
                    ("TOPPADDING", (0, 0), (-1, -1), 4),
                    ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
                ]
            )
        )
        story.append(doc_tbl)

    # Tasks
    pending = [t for t in tasks if t.get("status") == "pending"]
    if pending:
        story.append(Paragraph("Open Tasks", h1_style))
        for task in pending[:40]:
            text = (task.get("text") or "").strip()[:200]
            text = text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
            story.append(Paragraph(f"☐  {text}", bullet_style))

    doc.build(story)

    text = _format_knowledge_text(knowledge)
    title_out = f"Research Report — {work.get('title', 'Work')}"
    doc_id = _register_output(fpath, work_id, db, cfg, "pdf", title_out, text)
    logger.info("Generated PDF: %s → doc %s", fpath.name, doc_id)
    return fpath, doc_id


# ── PPTX ───────────────────────────────────────────────────────────────────────


def generate_pptx(work_id: str, db: OrivellumDB, cfg: OrivellumConfig) -> tuple[Path, str]:
    """Generate a PowerPoint deck from a Work; return (file_path, doc_id)."""
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    work = db.get_work(work_id)
    if not work:
        raise ValueError(f"Work {work_id!r} not found")

    knowledge = db.list_knowledge(work_id=work_id, limit=500)
    tasks = db.list_tasks(work_id=work_id)
    with db._lock:
        doc_rows = db._conn.execute(
            "SELECT title, kind FROM documents WHERE work_id=? ORDER BY created_at DESC LIMIT 30",
            (work_id,),
        ).fetchall()
    doc_list = [dict(r) for r in doc_rows]

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # ── Color palette ──
    DARK_BG = RGBColor(0x0F, 0x17, 0x2A)
    LIGHT_BG = RGBColor(0xF8, 0xFA, 0xFC)
    ACCENT = RGBColor(0x60, 0x81, 0xEB)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    GRAY = RGBColor(0x64, 0x74, 0x8B)
    DARK_TEXT = RGBColor(0x1E, 0x29, 0x3B)

    slide_layouts = prs.slide_layouts

    def _bg_slide(slide, color: RGBColor) -> None:
        """Fill a slide's background with a solid color."""
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = color

    def _add_text_box(
        slide,
        text: str,
        left,
        top,
        width,
        height,
        size: int = 14,
        bold: bool = False,
        color: RGBColor = WHITE,
        align=PP_ALIGN.LEFT,
        wrap: bool = True,
    ) -> None:
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = wrap
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color

    Inches(1)

    # ── Slide 1: Title ──
    slide1 = prs.slides.add_slide(slide_layouts[6])  # blank
    _bg_slide(slide1, DARK_BG)
    _add_text_box(
        slide1,
        work.get("title", "Research Report"),
        Inches(1),
        Inches(2.2),
        Inches(11.33),
        Inches(1.8),
        size=40,
        bold=True,
        color=WHITE,
        align=PP_ALIGN.CENTER,
    )
    desc = work.get("description") or ""
    if desc:
        _add_text_box(
            slide1,
            desc[:200],
            Inches(2),
            Inches(4.2),
            Inches(9.33),
            Inches(1),
            size=16,
            bold=False,
            color=ACCENT,
            align=PP_ALIGN.CENTER,
        )
    _add_text_box(
        slide1,
        f"Generated {datetime.now(UTC).strftime('%B %d, %Y')}",
        Inches(1),
        Inches(6.6),
        Inches(11.33),
        Inches(0.5),
        size=11,
        bold=False,
        color=GRAY,
        align=PP_ALIGN.CENTER,
    )

    # ── Slide 2: Overview ──
    slide2 = prs.slides.add_slide(slide_layouts[6])
    _bg_slide(slide2, LIGHT_BG)
    _add_text_box(
        slide2,
        "Overview",
        Inches(1),
        Inches(0.5),
        Inches(11.33),
        Inches(0.8),
        size=28,
        bold=True,
        color=DARK_TEXT,
    )
    stats = [
        (str(len(doc_list)), "Documents"),
        (str(len(knowledge)), "Knowledge Items"),
        (str(sum(1 for t in tasks if t.get("status") == "pending")), "Open Tasks"),
    ]
    for i, (val, lbl) in enumerate(stats):
        left = Inches(1 + i * 4)
        # Value
        v_box = slide2.shapes.add_textbox(left, Inches(2), Inches(3.5), Inches(1.5))
        v_tf = v_box.text_frame
        vp = v_tf.paragraphs[0]
        vp.alignment = PP_ALIGN.CENTER
        vr = vp.add_run()
        vr.text = val
        vr.font.size = Pt(54)
        vr.font.bold = True
        vr.font.color.rgb = ACCENT
        # Label
        l_box = slide2.shapes.add_textbox(left, Inches(3.6), Inches(3.5), Inches(0.6))
        l_tf = l_box.text_frame
        lp = l_tf.paragraphs[0]
        lp.alignment = PP_ALIGN.CENTER
        lr = lp.add_run()
        lr.text = lbl
        lr.font.size = Pt(14)
        lr.font.color.rgb = GRAY

    # ── Slides 3+: Knowledge by kind ──
    from collections import defaultdict

    by_kind: dict[str, list[dict]] = defaultdict(list)
    for item in knowledge:
        by_kind[item.get("kind", "note")].append(item)

    kind_order = ["fact", "claim", "note", "question", "action", "entity", "summary"]
    all_kinds = kind_order + [k for k in by_kind if k not in kind_order]

    for kind in all_kinds:
        items = by_kind.get(kind, [])
        if not items:
            continue
        # One slide per kind, up to 10 items
        slide = prs.slides.add_slide(slide_layouts[6])
        _bg_slide(slide, LIGHT_BG)
        _add_text_box(
            slide,
            f"{kind.title()}s",
            Inches(1),
            Inches(0.4),
            Inches(11.33),
            Inches(0.8),
            size=26,
            bold=True,
            color=DARK_TEXT,
        )
        # Bullet list
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1.4), Inches(11.33), Inches(5.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        first = True
        for item in items[:8]:
            text = (item.get("text") or "").strip()[:220]
            if not text:
                continue
            if first:
                p = tf.paragraphs[0]
                first = False
            else:
                p = tf.add_paragraph()
            p.space_before = Pt(4)
            run = p.add_run()
            run.text = f"• {text}"
            run.font.size = Pt(13)
            run.font.color.rgb = DARK_TEXT

    # ── Slide: Source Documents ──
    if doc_list:
        slide_d = prs.slides.add_slide(slide_layouts[6])
        _bg_slide(slide_d, DARK_BG)
        _add_text_box(
            slide_d,
            "Source Documents",
            Inches(1),
            Inches(0.4),
            Inches(11.33),
            Inches(0.8),
            size=26,
            bold=True,
            color=WHITE,
        )
        txBox = slide_d.shapes.add_textbox(Inches(1), Inches(1.4), Inches(11.33), Inches(5.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        first = True
        for d in doc_list[:12]:
            title = (d.get("title") or "")[:80]
            kind = d.get("kind", "")
            if first:
                p = tf.paragraphs[0]
                first = False
            else:
                p = tf.add_paragraph()
            p.space_before = Pt(3)
            run = p.add_run()
            run.text = f"• {title}" + (f"  [{kind}]" if kind else "")
            run.font.size = Pt(12)
            run.font.color.rgb = RGBColor(0xCB, 0xD5, 0xE1)

    # ── Slide: Open Tasks ──
    pending = [t for t in tasks if t.get("status") == "pending"]
    if pending:
        slide_t = prs.slides.add_slide(slide_layouts[6])
        _bg_slide(slide_t, DARK_BG)
        _add_text_box(
            slide_t,
            "Open Tasks",
            Inches(1),
            Inches(0.4),
            Inches(11.33),
            Inches(0.8),
            size=26,
            bold=True,
            color=WHITE,
        )
        txBox = slide_t.shapes.add_textbox(Inches(1), Inches(1.4), Inches(11.33), Inches(5.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        first = True
        for task in pending[:10]:
            text = (task.get("text") or "").strip()[:200]
            if first:
                p = tf.paragraphs[0]
                first = False
            else:
                p = tf.add_paragraph()
            run = p.add_run()
            run.text = f"☐  {text}"
            run.font.size = Pt(13)
            run.font.color.rgb = RGBColor(0xCB, 0xD5, 0xE1)

    # ── Save ──
    out_dir = _ensure_dir(cfg, work_id)
    slug = _slug(work.get("title", "work"))
    fname = f"{slug}_{_now_label()}.pptx"
    fpath = out_dir / fname
    prs.save(str(fpath))

    text = _format_knowledge_text(knowledge)
    title_out = f"Slide Deck — {work.get('title', 'Work')}"
    doc_id = _register_output(fpath, work_id, db, cfg, "pptx", title_out, text)
    logger.info("Generated PPTX: %s → doc %s", fpath.name, doc_id)
    return fpath, doc_id


# ── Bundle (ZIP) ───────────────────────────────────────────────────────────────


def bundle_files(
    file_paths: list[str],
    output_name: str,
    work_id: str,
    db: OrivellumDB,
    cfg: OrivellumConfig,
) -> tuple[Path, str]:
    """Zip a list of file paths into a single archive; return (zip_path, doc_id)."""
    out_dir = _ensure_dir(cfg, work_id)
    safe_name = _slug(output_name or "bundle")
    fname = f"{safe_name}_{_now_label()}.zip"
    fpath = out_dir / fname

    data_dir = Path(cfg.data_dir)
    # Restrict bundle inputs to files already generated for this specific Work.
    # This prevents an authenticated caller from packaging unrelated application
    # data (e.g. the database or another Work's outputs).
    work_gen_root = (data_dir / "outputs" / "generate" / work_id).resolve()
    included_text_parts: list[str] = []
    rejected: list[str] = []

    with zipfile.ZipFile(str(fpath), "w", zipfile.ZIP_DEFLATED) as zf:
        for raw_path in file_paths:
            p = Path(raw_path)
            # Normalise relative paths (callers pass relative-to-data_dir paths)
            if not p.is_absolute():
                p = data_dir / p
            try:
                target = p.resolve()
                # Hard scope: must be under data/outputs/generate/{work_id}/
                target.relative_to(work_gen_root)
                if not (target.exists() and target.is_file()):
                    rejected.append(raw_path)
                    logger.warning("Bundle: path not found — %s", raw_path)
                    continue
                zf.write(str(target), target.name)
                if target.suffix.lower() in (".txt", ".md"):
                    try:
                        included_text_parts.append(target.read_text(errors="replace")[:500])
                    except Exception:
                        pass
            except ValueError:
                # Path escapes the Work's generation directory
                rejected.append(raw_path)
                logger.warning("Bundle: rejected out-of-scope path — %s", raw_path)
            except Exception as exc:
                rejected.append(raw_path)
                logger.warning("Bundle: skipping %s — %s", raw_path, exc)

    if rejected:
        logger.info("Bundle: %d path(s) rejected (out-of-scope or missing)", len(rejected))

    text = "\n\n".join(included_text_parts) or f"Bundle containing {len(file_paths)} files."
    work = db.get_work(work_id)
    work_title = (work.get("title") if work else None) or "Work"
    title_out = f"Bundle — {work_title}"
    doc_id = _register_output(fpath, work_id, db, cfg, "bundle", title_out, text)
    logger.info("Generated ZIP bundle: %s → doc %s", fpath.name, doc_id)
    return fpath, doc_id


# ── Prompt-driven generation ───────────────────────────────────────────────────


def _build_docx_from_data(data: dict, out_path: Path) -> str:
    """Create a DOCX from LLM-structured JSON and return plain-text content."""
    from docx import Document as _Doc

    doc = _Doc()

    title = data.get("title", "Document")
    doc.add_heading(title, level=0)
    text_parts = [title]

    for section in data.get("sections", []):
        heading = section.get("heading", "")
        if heading:
            doc.add_heading(heading, level=1)
            text_parts.append(heading)
        for para in section.get("paragraphs", []):
            if str(para).strip():
                doc.add_paragraph(str(para))
                text_parts.append(str(para))

    doc.save(str(out_path))
    return "\n\n".join(text_parts)


def _build_pdf_from_data(data: dict, out_path: Path) -> str:
    """Create a PDF from LLM-structured JSON and return plain-text content."""
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.lib.units import cm
    from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer

    doc = SimpleDocTemplate(
        str(out_path),
        pagesize=A4,
        leftMargin=2.5 * cm,
        rightMargin=2.5 * cm,
        topMargin=2.5 * cm,
        bottomMargin=2.5 * cm,
    )
    styles = getSampleStyleSheet()
    story = []
    text_parts = []

    title = data.get("title", "Document")
    story.append(Paragraph(title, styles["Title"]))
    story.append(Spacer(1, 0.5 * cm))
    text_parts.append(title)

    for section in data.get("sections", []):
        heading = section.get("heading", "")
        if heading:
            story.append(Paragraph(heading, styles["Heading1"]))
            story.append(Spacer(1, 0.2 * cm))
            text_parts.append(heading)
        for para in section.get("paragraphs", []):
            if str(para).strip():
                story.append(
                    Paragraph(
                        str(para).replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;"),
                        styles["BodyText"],
                    )
                )
                story.append(Spacer(1, 0.15 * cm))
                text_parts.append(str(para))

    doc.build(story)
    return "\n\n".join(text_parts)


def _build_pptx_from_data(data: dict, out_path: Path) -> str:
    """Create a PPTX from LLM-structured JSON and return plain-text content."""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    text_parts = []
    title_str = str(data.get("title", "Presentation"))
    subtitle_str = str(data.get("subtitle", ""))

    # ── Title slide ──
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title_str
    try:
        slide.placeholders[1].text = subtitle_str
    except Exception:
        pass
    text_parts += [title_str, subtitle_str]

    # ── Content slides ──
    for s in data.get("slides", []):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide_title = str(s.get("title", ""))
        slide.shapes.title.text = slide_title
        text_parts.append(slide_title)

        try:
            tf = slide.placeholders[1].text_frame
            tf.clear()
            bullets = s.get("bullets", [])
            for i, bullet in enumerate(bullets):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = str(bullet)
                p.level = 0
                text_parts.append(str(bullet))
        except Exception:
            pass

        notes = str(s.get("speaker_notes", ""))
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    prs.save(str(out_path))
    return "\n\n".join(t for t in text_parts if t)


def _build_xlsx_from_data(data: dict, out_path: Path) -> str:
    """Create an XLSX from LLM-structured JSON and return plain-text content."""
    import openpyxl
    from openpyxl.styles import Alignment, Font, PatternFill

    wb = openpyxl.Workbook()
    default_ws = wb.active
    if default_ws:
        wb.remove(default_ws)

    text_parts = [str(data.get("title", "Spreadsheet"))]

    for sheet_data in data.get("sheets", []):
        name = str(sheet_data.get("name", "Sheet"))[:31]
        ws = wb.create_sheet(title=name)

        headers = sheet_data.get("headers", [])
        if headers:
            for col, h in enumerate(headers, 1):
                cell = ws.cell(row=1, column=col, value=str(h))
                cell.font = Font(bold=True, color="FFFFFF")
                cell.fill = PatternFill("solid", fgColor="1E293B")
                cell.alignment = Alignment(horizontal="center")
            text_parts.extend(str(h) for h in headers)

        for row_i, row in enumerate(sheet_data.get("rows", []), 2):
            for col_i, val in enumerate(row, 1):
                # Preserve numeric types; convert everything else to string
                if isinstance(val, (int, float)):
                    ws.cell(row=row_i, column=col_i, value=val)
                else:
                    ws.cell(row=row_i, column=col_i, value=str(val) if val is not None else "")

        summary = str(sheet_data.get("summary", ""))
        if summary:
            text_parts.append(summary)

    if not wb.sheetnames:
        wb.create_sheet("Sheet1")

    wb.save(str(out_path))
    return "\n".join(text_parts)


def generate_from_prompt(
    prompt: str,
    format: str,
    filename: str | None,
    work_id: str | None,
    db: OrivellumDB,
    cfg: OrivellumConfig,
) -> tuple[Path, str]:
    """Generate a document from a free-form text prompt using the LLM.

    Steps:
    1. Ask the LLM to produce structured JSON content for the requested format.
    2. Build the actual file with the appropriate library.
    3. Register it in the library (ARTIFACT tier) and return (path, doc_id).
    """
    import json
    import re

    from orivellum.capabilities.llm import llm_call

    fmt = format.lower().strip(".")
    if fmt not in ("docx", "pdf", "pptx", "xlsx"):
        raise ValueError(f"Unsupported format {fmt!r} — use docx, pdf, pptx, or xlsx")

    scope = work_id or "chat"
    out_dir = Path(cfg.data_dir) / "outputs" / "generate" / scope
    out_dir.mkdir(parents=True, exist_ok=True)

    stem = _slug(filename.rsplit(".", 1)[0] if filename else prompt[:60]) or "document"
    out_path = out_dir / f"{stem}_{_now_label()}.{fmt}"

    # ── Format-specific system prompt ──
    if fmt in ("docx", "pdf"):
        system = (
            "You are a professional writer. Generate rich, well-structured document content "
            "based on the user's request. Return ONLY valid JSON — no markdown fences, no prose "
            "outside the JSON object:\n"
            '{"title":"...","sections":[{"heading":"...","paragraphs":["...","..."]}]}'
        )
    elif fmt == "pptx":
        system = (
            "You are a presentation designer. Generate engaging PowerPoint content. "
            "Return ONLY valid JSON — no markdown fences, no prose outside the JSON object:\n"
            '{"title":"...","subtitle":"...","slides":[{"title":"...","bullets":["..."],'
            '"speaker_notes":"..."}]}\n'
            "Aim for 6-12 slides, 3-5 bullets each."
        )
    else:  # xlsx
        system = (
            "You are a data analyst. Generate spreadsheet data based on the user's request. "
            "Return ONLY valid JSON — no markdown fences, no prose outside the JSON object:\n"
            '{"title":"...","sheets":[{"name":"...","headers":["Col1","Col2"],'
            '"rows":[["val1","val2"]],"summary":"..."}]}'
        )

    result = llm_call(
        messages=[
            {"role": "system", "content": system},
            {"role": "user", "content": prompt},
        ],
        cfg=cfg,
        db=db,
        purpose=f"generate_{fmt}_prompt",
        timeout=90,
        max_tokens=4096,
        temperature=0.7,
    )

    if not result.ok:
        raise RuntimeError(f"LLM call failed: {result.error}")

    raw = (result.text or "").strip()
    # Strip any accidental markdown fences the model adds
    raw = re.sub(r"^```(?:json)?\s*\n?", "", raw, flags=re.MULTILINE)
    raw = re.sub(r"\n?```\s*$", "", raw, flags=re.MULTILINE)
    raw = raw.strip()

    try:
        data = json.loads(raw)
    except json.JSONDecodeError as exc:
        raise RuntimeError(f"LLM returned invalid JSON: {exc}\nPreview: {raw[:400]}")

    # ── Build file ──
    builders = {
        "docx": _build_docx_from_data,
        "pdf": _build_pdf_from_data,
        "pptx": _build_pptx_from_data,
        "xlsx": _build_xlsx_from_data,
    }
    text_content = builders[fmt](data, out_path)

    doc_title = str(data.get("title") or stem)
    doc_id = _register_output(out_path, work_id, db, cfg, fmt, doc_title, text_content)
    logger.info("generate_from_prompt: %s → %s (doc %s)", fmt, out_path.name, doc_id)
    return out_path, doc_id
