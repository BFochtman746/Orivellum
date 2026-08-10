"""File browser routes — /api/files/* and /api/upload, /api/download/*"""

from __future__ import annotations

import base64
import logging
import mimetypes
from pathlib import Path

from fastapi import APIRouter, Depends, HTTPException
from fastapi.responses import FileResponse
from pydantic import BaseModel

from orivellum.api._deps import get_config, require_auth

logger = logging.getLogger(__name__)

router = APIRouter(prefix="/api", dependencies=[Depends(require_auth)])


def _data_dir() -> Path:
    cfg = get_config()
    return Path(cfg.data_dir)


def _resolve_within_data_dir(rel_path: str) -> tuple[Path, Path]:
    """Resolve *rel_path* under the data dir, rejecting any escape.

    Uses ``Path.relative_to`` on fully resolved paths (never string-prefix
    comparison, which a sibling dir like ``data_backup/`` would bypass).
    Returns ``(data_dir_resolved, target_resolved)`` or raises 403.
    """
    data_dir = _data_dir().resolve()
    target = (data_dir / rel_path).resolve()
    try:
        target.relative_to(data_dir)
    except ValueError:
        raise HTTPException(403, "Path traversal not allowed")
    return data_dir, target


@router.get("/files")
def list_files(subdir: str = ""):
    data_dir, target = _resolve_within_data_dir(subdir)
    if not target.exists():
        return {"files": [], "dirs": [], "path": subdir}

    files = []
    dirs = []
    for item in sorted(target.iterdir()):
        if item.name.startswith("."):
            continue
        rel = str(item.relative_to(data_dir))
        if item.is_dir():
            dirs.append({"name": item.name, "path": rel})
        else:
            stat = item.stat()
            files.append(
                {
                    "name": item.name,
                    "path": rel,
                    "size_bytes": stat.st_size,
                    "mime": mimetypes.guess_type(item.name)[0] or "application/octet-stream",
                }
            )
    return {"files": files, "dirs": dirs, "path": subdir}


class UploadRequest(BaseModel):
    filename: str
    content_b64: str
    subdir: str = "intake"


@router.post("/upload")
def upload_file(body: UploadRequest):
    data_dir = _data_dir()
    try:
        data = base64.b64decode(body.content_b64, validate=True)
    except Exception:
        raise HTTPException(400, "content_b64 is not valid base64")

    name = Path(body.filename).name
    if not name or name.startswith("."):
        raise HTTPException(400, f"Bad filename: {body.filename!r}")

    _, target_dir = _resolve_within_data_dir(body.subdir)
    target_dir.mkdir(parents=True, exist_ok=True)

    dest = target_dir / name
    # Handle filename collision
    counter = 1
    while dest.exists():
        stem = Path(name).stem
        suffix = Path(name).suffix
        dest = target_dir / f"{stem}_{counter}{suffix}"
        counter += 1

    dest.write_bytes(data)
    return {
        "ok": True,
        "filename": dest.name,
        "path": str(dest.relative_to(data_dir)),
        "size_bytes": len(data),
    }


class ExtractFileRequest(BaseModel):
    filename: str
    content_b64: str


@router.post("/extract-file")
async def extract_file_for_chat(body: ExtractFileRequest):
    """Extract readable text from a base64-encoded document for chat context injection.

    Supports: .pdf, .docx, .xlsx, .csv, .txt, .md and plain-text variants.
    The extracted text is capped at 12 000 characters to stay within the model context budget.
    """
    import io

    _MAX_CHARS = 12_000

    try:
        raw = base64.b64decode(body.content_b64, validate=True)
    except Exception:
        raise HTTPException(400, "content_b64 is not valid base64")

    from pathlib import Path as _P

    ext = _P(body.filename).suffix.lower()
    text = ""

    try:
        if ext in (".txt", ".md", ".csv", ".json", ".yaml", ".yml", ".xml", ".html", ".htm"):
            text = raw.decode("utf-8", errors="replace")

        elif ext == ".pdf":
            try:
                import fitz

                doc = fitz.open(stream=raw, filetype="pdf")
                parts = [page.get_text().strip() for page in doc if page.get_text().strip()]
                text = "\n\n".join(parts)
                doc.close()
            except Exception:
                try:
                    import pdfminer.high_level as _pml

                    text = _pml.extract_text(io.BytesIO(raw)) or ""
                except Exception:
                    logger.exception("PDF extraction failed for %s", body.filename)
                    text = "[PDF extraction failed — see server logs]"

        elif ext in (".docx",):
            try:
                from docx import Document as _Doc

                doc = _Doc(io.BytesIO(raw))
                paras = [p.text for p in doc.paragraphs if p.text.strip()]
                for table in doc.tables:
                    for row in table.rows:
                        paras.append("\t".join(c.text.strip() for c in row.cells))
                text = "\n".join(paras)
            except Exception:
                logger.exception("DOCX extraction failed for %s", body.filename)
                text = "[DOCX extraction failed — see server logs]"

        elif ext in (".xlsx", ".xls"):
            try:
                import openpyxl

                wb = openpyxl.load_workbook(io.BytesIO(raw), read_only=True, data_only=True)
                rows_out: list[str] = []
                for ws in wb.worksheets:
                    rows_out.append(f"## Sheet: {ws.title}")
                    for i, row in enumerate(ws.iter_rows(values_only=True)):
                        if i >= 500:
                            rows_out.append("[... row limit reached (500) ...]")
                            break
                        rows_out.append("\t".join("" if c is None else str(c) for c in row))
                text = "\n".join(rows_out)
            except Exception:
                logger.exception("Spreadsheet extraction failed for %s", body.filename)
                text = "[Spreadsheet extraction failed — see server logs]"

        elif ext in (".pptx",):
            try:
                from pptx import Presentation as _Prs

                prs = _Prs(io.BytesIO(raw))
                slides_text: list[str] = []
                for i, slide in enumerate(prs.slides, 1):
                    slide_parts: list[str] = [f"--- Slide {i} ---"]
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            slide_parts.extend(
                                p.text for p in shape.text_frame.paragraphs if p.text.strip()
                            )
                    slides_text.append("\n".join(slide_parts))
                text = "\n\n".join(slides_text)
            except Exception:
                logger.exception("PPTX extraction failed for %s", body.filename)
                text = "[PPTX extraction failed — see server logs]"

        else:
            text = (
                f"[File type '{ext}' is not supported for text extraction. "
                "Please use PDF, DOCX, PPTX, XLSX, CSV, or plain text.]"
            )
    except Exception:
        # FA-03: never leak raw exception text to the client — the marker is
        # embedded in extracted_text (chat context), so keep it generic.
        logger.exception("extract-file failed for %s", body.filename)
        text = "[Extraction failed — see server logs]"

    if len(text) > _MAX_CHARS:
        text = text[:_MAX_CHARS] + f"\n\n[... content truncated to {_MAX_CHARS:,} characters ...]"

    return {
        "ok": True,
        "filename": body.filename,
        "extracted_text": text,
        "char_count": len(text),
    }


# FA-02: only user-content subtrees are servable. Everything else under the
# data dir — SQLite DBs (+ -wal/-shm), backups, the mail token vault,
# api_key.txt — must never leave the server via this endpoint.
_DOWNLOAD_ALLOWED_SUBDIRS = frozenset({"library", "outputs", "intake"})
_DOWNLOAD_DENIED_SUFFIXES = (".db", ".db-wal", ".db-shm", ".sqlite", ".sqlite3")
_DOWNLOAD_DENIED_NAMES = frozenset({"api_key.txt"})


@router.get("/download/{path:path}")
def download_file(path: str):
    data_dir, target = _resolve_within_data_dir(path)
    rel = target.relative_to(data_dir)
    if not rel.parts or rel.parts[0] not in _DOWNLOAD_ALLOWED_SUBDIRS:
        raise HTTPException(
            403,
            "Downloads are limited to: " + ", ".join(sorted(_DOWNLOAD_ALLOWED_SUBDIRS)),
        )
    lower = target.name.lower()
    if lower in _DOWNLOAD_DENIED_NAMES or lower.endswith(_DOWNLOAD_DENIED_SUFFIXES):
        raise HTTPException(403, "This file type is not downloadable")
    if not target.exists() or not target.is_file():
        raise HTTPException(404, f"File not found: {path}")
    return FileResponse(str(target), filename=target.name)
